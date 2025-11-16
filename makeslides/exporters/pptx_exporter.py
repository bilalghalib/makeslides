#!/usr/bin/env python3
"""
PowerPoint (PPTX) exporter using python-pptx.

Creates professional PowerPoint presentations with full control over layouts,
images can be embedded directly (no external hosting needed).
"""
import io
import logging
import requests
from pathlib import Path
from typing import List, Dict, Any, Optional
from urllib.parse import urlparse

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
except ImportError:
    raise ImportError("python-pptx is required for PPTX export. Install with: pip install python-pptx")

from .base import BaseExporter

logger = logging.getLogger(__name__)


class PPTXExporter(BaseExporter):
    """Export presentations to PowerPoint (PPTX) format."""

    # Standard PowerPoint dimensions (16:9)
    SLIDE_WIDTH = Inches(10)
    SLIDE_HEIGHT = Inches(7.5)

    # Color scheme (customizable)
    PRIMARY_COLOR = RGBColor(0, 120, 212)  # Blue
    SECONDARY_COLOR = RGBColor(51, 51, 51)  # Dark gray
    ACCENT_COLOR = RGBColor(255, 185, 0)  # Orange

    def __init__(self, slides_data: List[Dict[str, Any]], output_path: Optional[Path] = None,
                 theme: str = "modern"):
        """
        Initialize PPTX exporter.

        Args:
            slides_data: List of slide dictionaries
            output_path: Path for output PPTX file
            theme: Theme name (modern, classic, minimal)
        """
        super().__init__(slides_data, output_path)
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT
        self.theme = theme

    def export(self) -> Path:
        """
        Export presentation to PPTX file.

        Returns:
            Path to the generated PPTX file
        """
        if not self.validate_slides():
            raise ValueError("Invalid slides data")

        logger.info(f"Exporting {len(self.slides_data)} slides to PPTX...")

        for i, slide_data in enumerate(self.slides_data, 1):
            logger.info(f"Creating slide {i}/{len(self.slides_data)}: {slide_data.get('title', 'Untitled')}")
            self._create_slide(slide_data)

        # Determine output path
        if not self.output_path:
            title_slug = self.metadata['title'].lower().replace(' ', '_')
            self.output_path = Path(f"{title_slug}.pptx")

        # Save presentation
        self.prs.save(str(self.output_path))
        logger.info(f"✅ PPTX presentation saved to: {self.output_path}")

        return self.output_path

    def _create_slide(self, slide_data: Dict[str, Any]):
        """Create a slide based on its layout type."""
        layout_type = self._get_layout_type(slide_data)

        # Map layout types to creation methods
        layout_methods = {
            'title': self._create_title_slide,
            'section': self._create_section_slide,
            'content': self._create_content_slide,
            'two_columns': self._create_two_column_slide,
            'quote': self._create_quote_slide,
            'main_point': self._create_main_point_slide,
            'big_number': self._create_big_number_slide,
            'caption': self._create_caption_slide,
            'blank': self._create_blank_slide
        }

        method = layout_methods.get(layout_type, self._create_content_slide)
        method(slide_data)

    def _create_title_slide(self, slide_data: Dict[str, Any]):
        """Create a title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = slide_data.get('title', '')
        subtitle.text = slide_data.get('content', '')

        # Style title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.PRIMARY_COLOR

        # Style subtitle
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.SECONDARY_COLOR

    def _create_section_slide(self, slide_data: Dict[str, Any]):
        """Create a section header slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])  # Section header layout

        title = slide.shapes.title
        title.text = slide_data.get('title', '')

        # Style as large, centered text
        title.text_frame.paragraphs[0].font.size = Pt(54)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.PRIMARY_COLOR
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add colored bar at bottom for visual interest
        left = Inches(0)
        top = Inches(6.8)
        width = Inches(10)
        height = Inches(0.7)

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.ACCENT_COLOR
        bar.line.fill.background()

    def _create_content_slide(self, slide_data: Dict[str, Any]):
        """Create a standard content slide with bullet points."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])  # Title and Content layout

        title = slide.shapes.title
        title.text = slide_data.get('title', '')

        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()

        # Add bullet points
        bullets = self._format_bullet_points(slide_data.get('content', ''))

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = bullet
            p.font.size = Pt(18)
            p.level = 0

        # Add image if present
        self._add_image_to_slide(slide, slide_data)

    def _create_two_column_slide(self, slide_data: Dict[str, Any]):
        """Create a two-column layout slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])

        title = slide.shapes.title
        title.text = slide_data.get('title', '')

        # Remove default content placeholder
        for shape in slide.shapes:
            if shape.has_text_frame and shape != title:
                sp = shape.element
                sp.getparent().remove(sp)

        # Create two columns
        left_col_left = Inches(0.5)
        left_col_top = Inches(1.5)
        left_col_width = Inches(4.5)
        left_col_height = Inches(5)

        right_col_left = Inches(5.2)
        right_col_top = Inches(1.5)
        right_col_width = Inches(4.5)
        right_col_height = Inches(5)

        # Left column
        left_box = slide.shapes.add_textbox(left_col_left, left_col_top, left_col_width, left_col_height)
        left_frame = left_box.text_frame
        left_frame.word_wrap = True

        # Right column
        right_box = slide.shapes.add_textbox(right_col_left, right_col_top, right_col_width, right_col_height)
        right_frame = right_box.text_frame
        right_frame.word_wrap = True

        # Split content or use provided columns
        content = slide_data.get('content', '')

        if '|' in content:
            left_content, right_content = content.split('|', 1)
        else:
            # Split bullet points in half
            bullets = self._format_bullet_points(content)
            mid = len(bullets) // 2
            left_content = '\n'.join(bullets[:mid])
            right_content = '\n'.join(bullets[mid:])

        # Add left content
        for bullet in left_content.strip().split('\n'):
            if bullet.strip():
                p = left_frame.add_paragraph()
                p.text = bullet.strip().lstrip('*-• ').strip()
                p.font.size = Pt(16)

        # Add right content or image
        image_url = slide_data.get('image_url')
        if image_url and not 'diagram' in image_url.lower():
            # Add image to right column
            self._add_image_to_slide(slide, slide_data, left=right_col_left, top=right_col_top,
                                    width=right_col_width, height=right_col_height)
        else:
            # Add text to right column
            for bullet in right_content.strip().split('\n'):
                if bullet.strip():
                    p = right_frame.add_paragraph()
                    p.text = bullet.strip().lstrip('*-• ').strip()
                    p.font.size = Pt(16)

    def _create_quote_slide(self, slide_data: Dict[str, Any]):
        """Create a quote slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Add quote text
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(2)

        quote_box = slide.shapes.add_textbox(left, top, width, height)
        quote_frame = quote_box.text_frame
        quote_frame.word_wrap = True

        p = quote_frame.paragraphs[0]
        p.text = f'"{slide_data.get("content", "")}"'
        p.font.size = Pt(32)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = self.PRIMARY_COLOR

        # Add attribution
        attr_top = Inches(5)
        attr_box = slide.shapes.add_textbox(left, attr_top, width, Inches(1))
        attr_frame = attr_box.text_frame

        attr_p = attr_frame.paragraphs[0]
        attr_p.text = f"— {slide_data.get('title', '')}"
        attr_p.font.size = Pt(18)
        attr_p.alignment = PP_ALIGN.CENTER
        attr_p.font.color.rgb = self.SECONDARY_COLOR

    def _create_main_point_slide(self, slide_data: Dict[str, Any]):
        """Create a slide emphasizing a main point."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Large centered text
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(9)
        height = Inches(3)

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        p.text = slide_data.get('title', '')
        p.font.size = Pt(60)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = self.PRIMARY_COLOR

        # Add supporting text if present
        content = slide_data.get('content', '')
        if content:
            content_top = Inches(5.5)
            content_box = slide.shapes.add_textbox(left, content_top, width, Inches(1.5))
            content_frame = content_box.text_frame

            cp = content_frame.paragraphs[0]
            cp.text = content
            cp.font.size = Pt(20)
            cp.alignment = PP_ALIGN.CENTER
            cp.font.color.rgb = self.SECONDARY_COLOR

    def _create_big_number_slide(self, slide_data: Dict[str, Any]):
        """Create a slide for displaying statistics or big numbers."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Extract number from title (if present)
        title = slide_data.get('title', '')

        # Huge number
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(2.5)

        num_box = slide.shapes.add_textbox(left, top, width, height)
        num_frame = num_box.text_frame
        num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = num_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(88)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = self.ACCENT_COLOR

        # Description
        desc_top = Inches(4.5)
        desc_box = slide.shapes.add_textbox(left, desc_top, width, Inches(2))
        desc_frame = desc_box.text_frame

        dp = desc_frame.paragraphs[0]
        dp.text = slide_data.get('content', '')
        dp.font.size = Pt(24)
        dp.alignment = PP_ALIGN.CENTER
        dp.font.color.rgb = self.SECONDARY_COLOR

    def _create_caption_slide(self, slide_data: Dict[str, Any]):
        """Create a slide with image and caption."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Add image (larger)
        self._add_image_to_slide(slide, slide_data, left=Inches(1), top=Inches(0.5),
                                width=Inches(8), height=Inches(5.5))

        # Add caption
        caption_top = Inches(6.2)
        caption_box = slide.shapes.add_textbox(Inches(1), caption_top, Inches(8), Inches(1))
        caption_frame = caption_box.text_frame

        p = caption_frame.paragraphs[0]
        p.text = slide_data.get('title', '')
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = self.SECONDARY_COLOR

    def _create_blank_slide(self, slide_data: Dict[str, Any]):
        """Create a blank slide with background image."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Full-screen image
        self._add_image_to_slide(slide, slide_data, left=Inches(0), top=Inches(0),
                                width=self.SLIDE_WIDTH, height=self.SLIDE_HEIGHT)

    def _add_image_to_slide(self, slide, slide_data: Dict[str, Any],
                           left: Optional[Inches] = None, top: Optional[Inches] = None,
                           width: Optional[Inches] = None, height: Optional[Inches] = None):
        """
        Add an image to a slide.

        Images can be from URLs or local paths. They are embedded in the PPTX file.
        """
        image_url = slide_data.get('image_url')

        if not image_url:
            return

        # Default positioning (bottom right corner)
        if left is None:
            left = Inches(6)
        if top is None:
            top = Inches(4)
        if width is None:
            width = Inches(3.5)
        if height is None:
            height = Inches(2.5)

        try:
            # Check if it's a URL or local path
            if image_url.startswith('http://') or image_url.startswith('https://'):
                # Download image
                logger.info(f"Downloading image from {image_url}")
                response = requests.get(image_url, timeout=10)
                response.raise_for_status()
                image_stream = io.BytesIO(response.content)

                slide.shapes.add_picture(image_stream, left, top, width=width, height=height)

            else:
                # Local file
                if Path(image_url).exists():
                    slide.shapes.add_picture(str(image_url), left, top, width=width, height=height)
                else:
                    logger.warning(f"Image not found: {image_url}")

        except Exception as e:
            logger.error(f"Failed to add image {image_url}: {e}")


def export_to_pptx(slides_data: List[Dict[str, Any]], output_path: Optional[Path] = None) -> Path:
    """
    Convenience function to export slides to PPTX.

    Args:
        slides_data: List of slide dictionaries
        output_path: Output file path

    Returns:
        Path to generated PPTX file
    """
    exporter = PPTXExporter(slides_data, output_path)
    return exporter.export()
